"""Catch the slide faults a renderer would show, on a machine with no renderer.

The proper check is to convert the deck to images and look at them. There is
no LibreOffice on this machine, so that is not available — and "I could not
render it" is not the same as "it is fine". This measures what can be
measured without a rasteriser:

  * anything hanging off the slide, or inside the half-inch margin
  * text boxes that overlap each other
  * text that cannot fit the box it is in, estimated from the glyph budget

The last one is an ESTIMATE and says so. It uses average character widths per
point size, so it is wrong at the edges — a line of capitals is wider than the
average, a column of "i" narrower. It is tuned to complain early: a box at 85%
full is reported, because the cost of looking at a box that turned out fine is
nothing and the cost of shipping a cut-off sentence is a slide somebody reads
out in a pitch.
"""
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE

EMU = 914400.0
MARGIN = 0.5           # inches the skill asks content to stay inside
FULL = 0.85            # report a box estimated fuller than this


def inches(v):
    return (v or 0) / EMU


def text_of(shape):
    if not shape.has_text_frame:
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs)


def sizes(shape):
    """Every run's point size, defaulting to 18 where the run does not say."""
    out = []
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            out.append((r.text, (r.font.size.pt if r.font.size else 18),
                        bool(r.font.bold)))
    return out


def fill_ratio(shape):
    """Roughly how full the box is: 0.5 is half, >1 will not fit.

    Character width is taken as 0.5 of the point size, which is close for a
    humanist sans at mixed case. Line height is 1.22 of the point size unless
    the paragraph sets its own.
    """
    w = inches(shape.width) * 72.0        # points
    h = inches(shape.height) * 72.0
    if w <= 0 or h <= 0:
        return 0.0
    # pptxgenjs default inset is 0.05" each side unless margin:0 was set.
    w -= 4
    if w <= 0:
        return 9.9
    lines = 0.0
    tallest = 0.0
    for para in shape.text_frame.paragraphs:
        runs = [(r.text, (r.font.size.pt if r.font.size else 18),
                 bool(r.font.bold)) for r in para.runs]
        if not runs:
            lines += 1
            continue
        pt = max(s for _, s, _ in runs) or 18
        tallest = max(tallest, pt)
        width = sum(len(t) * s * (0.53 if b else 0.5) for t, s, b in runs)
        lines += max(1.0, -(-width // w) if w else 1.0)
    if not tallest:
        return 0.0
    return (lines * tallest * 1.22) / h


def main(path):
    prs = Presentation(path)
    W = inches(prs.slide_width)
    H = inches(prs.slide_height)
    print(f"{path}: {len(prs.slides.__iter__.__self__._sldIdLst)} slides, "
          f"{W:.2f} x {H:.2f} in")
    bad = 0

    for i, slide in enumerate(prs.slides, 1):
        notes = []
        boxes = []
        panels = []   # filled shapes carrying no text of their own
        for sh in slide.shapes:
            x, y = inches(sh.left), inches(sh.top)
            w, h = inches(sh.width), inches(sh.height)
            t = text_of(sh).strip()

            if x < -0.01 or y < -0.01 or x + w > W + 0.01 or y + h > H + 0.01:
                # A deliberate bleed shape is fine; text off the slide is not.
                if t:
                    notes.append(f"OFF-SLIDE text at ({x:.2f},{y:.2f}) "
                                 f"{w:.2f}x{h:.2f}: {t[:44]!r}")
            elif t and (x < MARGIN - 0.01 or y < MARGIN - 0.01
                        or x + w > W - MARGIN + 0.01
                        or y + h > H - MARGIN + 0.01):
                notes.append(f"inside the {MARGIN}in margin "
                             f"({x:.2f},{y:.2f} {w:.2f}x{h:.2f}): {t[:38]!r}")

            # Rectangles only. The round blobs on the title and closing
            # slides are backdrop, and text is MEANT to lie across them —
            # counting those as panels reported every cover as broken.
            if not t and w > 0.6 and h > 0.5 and not (w > 12.0 and h > 6.5):
                try:
                    kind = sh.auto_shape_type
                except Exception:
                    kind = None
                if kind in (MSO_SHAPE.ROUNDED_RECTANGLE, MSO_SHAPE.RECTANGLE):
                    panels.append((x, y, w, h))

            if t:
                r = fill_ratio(sh)
                if r > FULL:
                    notes.append(f"text may not fit (~{r*100:.0f}% of "
                                 f"{w:.2f}x{h:.2f}): {t[:46]!r}")
                boxes.append((x, y, w, h, t))

        # Text half-on a panel it does not belong to.
        #
        # Comparing text against text alone missed the fault this was written
        # for: a footer running under the bottom edge of a card. The card
        # holds no text of its own, so it was invisible to the check, and the
        # slide shipped with a sentence sliced by a rectangle.
        #
        # Containment is the test, not intersection — a caption INSIDE a card
        # overlaps it completely and is exactly right. What is wrong is text
        # that is partly on and partly off.
        for x, y, w, h, t in boxes:
            for px, py, pw, ph in panels:
                ox = min(x + w, px + pw) - max(x, px)
                oy = min(y + h, py + ph) - max(y, py)
                # Both directions must be a real overlap. A title grazing
                # a panel edge by four hundredths of an inch is the
                # rounding in the drawing, not a collision.
                if ox <= 0.09 or oy <= 0.09:
                    continue
                inside = (x >= px - 0.02 and y >= py - 0.02
                          and x + w <= px + pw + 0.02
                          and y + h <= py + ph + 0.02)
                if not inside and ox * oy > 0.02:
                    notes.append(f"text half-on a panel ({ox:.2f}x{oy:.2f}in): "
                                 f"{t[:44]!r}")

        # Text boxes sitting on top of each other.
        for a in range(len(boxes)):
            for b in range(a + 1, len(boxes)):
                ax, ay, aw, ah, at = boxes[a]
                bx, by, bw, bh, bt = boxes[b]
                ox = min(ax + aw, bx + bw) - max(ax, bx)
                oy = min(ay + ah, by + bh) - max(ay, by)
                if ox > 0.06 and oy > 0.06:
                    notes.append(f"text overlaps text ({ox:.2f}x{oy:.2f}in): "
                                 f"{at[:26]!r} / {bt[:26]!r}")

        if notes:
            bad += 1
            print(f"\nslide {i}")
            for n in notes:
                print("  " + n)

    print(f"\n{bad} slide(s) with something to look at")
    return 1 if bad else 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1]))
